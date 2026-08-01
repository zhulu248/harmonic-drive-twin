#!/usr/bin/env python3
"""Schaeffler-branded 3-slide deck: harmonic-drive-twin package anatomy.

Condenses docs/PACKAGE_ANATOMY.md into three slides.
Layout and brand helpers adapted from the NVIDIA weekly deck generator
(flex_base/nvidia_weekly_2026-07-23/build_deck.py).
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- brand constants ----
SCHAEFFLER_GREEN = RGBColor(0, 137, 61)
GREEN_LIGHT = RGBColor(230, 244, 236)
CARBON_GRAY = RGBColor(100, 100, 100)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(230, 230, 230)
FONT_NAME = "Arial"
MONO_NAME = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_LEFT = Inches(0.524)
CATEGORY_TOP = Inches(0.52)
TITLE_TOP = Inches(0.787)
TITLE_WIDTH = Inches(9.764)
GREEN_LINE_TOP = Inches(1.34)
CONTENT_TOP = Inches(1.47)
FOOTER_TOP = Inches(7.14)
LOGO_LEFT = Inches(11.037)
LOGO_TOP = Inches(0.522)
LOGO_W = Inches(1.772)
LOGO_H = Inches(0.198)
FULL_W = Inches(12.285)

BRAND = "/home/user01/Desktop/misc/branding/content_sheet"
LOGO_PATH = f"{BRAND}/logos/schaeffler_wordmark_green.png"
CLAIM_PATH = f"{BRAND}/logos/we_pioneer_motion_claim.png"

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "Harmonic_Drive_Twin_Package_Anatomy.pptx")

FOOTER_DATE = "July 2026"
FOOTER_DOC = "harmonic-drive-twin  |  Package anatomy  |  Schaeffler R&D"


# ---- helpers ----
def set_font(run, size=Pt(14), bold=False, color=CARBON_GRAY, name=FONT_NAME):
    run.font.size, run.font.bold, run.font.color.rgb, run.font.name = size, bold, color, name


def add_textbox(slide, left, top, width, height, text, size=Pt(14), bold=False,
                color=CARBON_GRAY, alignment=PP_ALIGN.LEFT, name=FONT_NAME):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    r = p.add_run()
    r.text = text
    set_font(r, size=size, bold=bold, color=color, name=name)
    return tb


def _no_shadow(shape):
    spPr = shape._element.spPr
    for s in spPr.findall(qn('a:effectLst')):
        spPr.remove(s)


def add_logo(slide):
    slide.shapes.add_picture(LOGO_PATH, LOGO_LEFT, LOGO_TOP, LOGO_W, LOGO_H)


def add_green_line(slide):
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), GREEN_LINE_TOP, SLIDE_W, Pt(1.5))
    ln.fill.solid(); ln.fill.fore_color.rgb = SCHAEFFLER_GREEN; ln.line.fill.background()
    _no_shadow(ln)


def add_marker(slide, left, top, height, width=Inches(0.12)):
    m = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    m.fill.solid(); m.fill.fore_color.rgb = SCHAEFFLER_GREEN; m.line.fill.background()
    _no_shadow(m)
    return m


def header(slide, category, title):
    add_textbox(slide, MARGIN_LEFT, CATEGORY_TOP, TITLE_WIDTH, Inches(0.13),
                "  ".join(category.upper()), size=Pt(8), color=CARBON_GRAY)
    add_textbox(slide, MARGIN_LEFT, TITLE_TOP, TITLE_WIDTH, Inches(0.4),
                title, size=Pt(20), bold=True, color=CARBON_GRAY)
    add_logo(slide)
    add_green_line(slide)


def footer(slide, n):
    add_textbox(slide, MARGIN_LEFT, FOOTER_TOP, Inches(0.9), Inches(0.14), FOOTER_DATE, size=Pt(7))
    add_textbox(slide, Inches(1.6), FOOTER_TOP, Inches(7.5), Inches(0.14), FOOTER_DOC, size=Pt(7))
    add_textbox(slide, Inches(9.0), FOOTER_TOP, Inches(3.5), Inches(0.14),
                "Copyright Schaeffler AG, 2026 – All rights reserved.  |", size=Pt(7),
                alignment=PP_ALIGN.RIGHT)
    add_textbox(slide, Inches(12.54), FOOTER_TOP, Inches(0.27), Inches(0.14), str(n),
                size=Pt(7), bold=True, alignment=PP_ALIGN.RIGHT)


def bullets(slide, left, top, width, height, items, size=Pt(15), spacing=Pt(10),
            color=CARBON_GRAY):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = "•  " + it
        set_font(r, size=size, color=color)
    return tb


def brand_table(slide, left, top, width, col_w, rows, header_row,
                head_size=Pt(11), body_size=Pt(10.5), row_h=Inches(0.34)):
    """Green-header table with alternating light-green rows.

    ``rows`` is a list of row tuples.  A cell may be either a plain string or a
    ``(headline, rest)`` tuple, which renders as bold-headline + normal text —
    the format used throughout PACKAGE_ANATOMY.md.
    """
    n_rows, n_cols = len(rows) + 1, len(col_w)
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width,
                                   row_h * n_rows)
    tbl = shape.table

    # strip PowerPoint's default blue table style / banding
    tblPr = tbl._tbl.tblPr
    for e in tblPr.findall(qn('a:tableStyleId')):
        tblPr.remove(e)
    tblPr.set('firstRow', '0')
    tblPr.set('bandRow', '0')

    for j, w in enumerate(col_w):
        tbl.columns[j].width = w
    tbl.rows[0].height = row_h

    def _fill(cell, rgb):
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb

    def _write(cell, content, size, bold_all=False, color=CARBON_GRAY,
               mono=False):
        cell.margin_left = Inches(0.09)
        cell.margin_right = Inches(0.07)
        cell.margin_top = Inches(0.03)
        cell.margin_bottom = Inches(0.03)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        if isinstance(content, tuple):
            head, rest = content
            r1 = p.add_run(); r1.text = head + " "
            set_font(r1, size=size, bold=True, color=color)
            r2 = p.add_run(); r2.text = rest
            set_font(r2, size=size, bold=False, color=color)
        else:
            r = p.add_run(); r.text = content
            set_font(r, size=size, bold=bold_all, color=color,
                     name=MONO_NAME if mono else FONT_NAME)

    for j, txt in enumerate(header_row):
        c = tbl.cell(0, j)
        _fill(c, SCHAEFFLER_GREEN)
        _write(c, txt, head_size, bold_all=True, color=WHITE)

    for i, row in enumerate(rows):
        tbl.rows[i + 1].height = row_h
        for j, txt in enumerate(row):
            c = tbl.cell(i + 1, j)
            _fill(c, GREEN_LIGHT if i % 2 == 0 else WHITE)
            _write(c, txt, body_size, mono=(j == 0))
    return shape


def keybox(slide, left, top, width, height, title, text,
           title_size=Pt(12), body_size=Pt(11.5)):
    """Green-tinted callout box with a green marker bar — the 'key result' look."""
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.solid(); box.fill.fore_color.rgb = GREEN_LIGHT
    box.line.color.rgb = SCHAEFFLER_GREEN; box.line.width = Pt(0.75)
    _no_shadow(box)
    box.text_frame.text = ""
    add_marker(slide, left, top, height, width=Inches(0.07))
    add_textbox(slide, left + Inches(0.2), top + Inches(0.09),
                width - Inches(0.35), Inches(0.24), title,
                size=title_size, bold=True, color=SCHAEFFLER_GREEN)
    tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.38),
                                  width - Inches(0.35), height - Inches(0.5))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    set_font(r, size=body_size, color=CARBON_GRAY)
    return box


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---- slides ----
def slide_overview(prs):
    """Slide 1 — what ships, and the four layers."""
    s = blank(prs)
    header(s, "Package anatomy – 1 / 3", "What ships, and how it is layered")

    add_textbox(s, MARGIN_LEFT, CONTENT_TOP + Inches(0.06), Inches(12.2), Inches(0.34),
                "A physics calculator for a strain-wave gearbox: give it an operating point, "
                "it returns the motor torque you need, the efficiency, where every watt went, "
                "and how much the drive twists.",
                size=Pt(13), color=CARBON_GRAY)

    # left: the four layers, bottom-up
    add_textbox(s, MARGIN_LEFT, Inches(2.06), Inches(5.6), Inches(0.3),
                "Four layers – each knows only the ones below it",
                size=Pt(13), bold=True, color=SCHAEFFLER_GREEN)

    layers = [
        ("4", "__init__.py", "the front door", "curates the 13-name public API"),
        ("3", "twin.py", "the engine", "does the work – solve() and 11 more questions"),
        ("2", "params.py", "the dials", "the only thing a user should edit"),
        ("1", "flex_common/", "the physics foundation", "the equations, 7 modules"),
    ]
    y = Inches(2.48)
    for num, mod, role, note in layers:
        add_marker(s, MARGIN_LEFT, y, Inches(0.62), width=Inches(0.09))
        add_textbox(s, MARGIN_LEFT + Inches(0.2), y - Inches(0.02), Inches(0.3), Inches(0.26),
                    num, size=Pt(13), bold=True, color=SCHAEFFLER_GREEN)
        add_textbox(s, MARGIN_LEFT + Inches(0.52), y - Inches(0.02), Inches(1.9), Inches(0.26),
                    mod, size=Pt(12), bold=True, color=CARBON_GRAY, name=MONO_NAME)
        add_textbox(s, MARGIN_LEFT + Inches(2.42), y - Inches(0.02), Inches(3.1), Inches(0.26),
                    role, size=Pt(12), bold=True, color=SCHAEFFLER_GREEN)
        add_textbox(s, MARGIN_LEFT + Inches(0.52), y + Inches(0.24), Inches(5.0), Inches(0.26),
                    note, size=Pt(11), color=CARBON_GRAY)
        y += Inches(0.78)

    # right: what a user actually opens
    x = Inches(6.6)
    add_textbox(s, x, Inches(2.06), Inches(5.9), Inches(0.3),
                "What is in the box", size=Pt(13), bold=True, color=SCHAEFFLER_GREEN)
    brand_table(
        s, x, Inches(2.44), Inches(6.2),
        [Inches(1.72), Inches(0.82), Inches(2.28), Inches(1.38)],
        [
            ("README.md", "249", "shop window", "yes, first"),
            ("__init__.py", "129", "front door / public API", "indirectly"),
            ("twin.py", "787", "the engine", "if curious"),
            ("params.py", "425", "the dials", "yes, to adapt"),
            ("catalogue.py", "250", "reality check", "to validate"),
            ("flex_common/", "1 356", "physics foundation", "rarely"),
            ("examples/", "503", "on-ramps", "yes, to start"),
            ("tests/", "183", "the proof", "to verify"),
        ],
        ["File", "Lines", "Role", "User opens it?"],
        body_size=Pt(10), row_h=Inches(0.335),
    )
    add_textbox(s, x, Inches(5.62), Inches(6.2), Inches(0.5),
                "2 947 lines of Python in the model itself. NumPy + SciPy only – "
                "no compiled extensions, no companion repo, no license server.",
                size=Pt(10.5), color=CARBON_GRAY)
    footer(s, 1)


def slide_foundation(prs):
    """Slide 2 — the flex_common module table (the centrepiece)."""
    s = blank(prs)
    header(s, "Package anatomy – 2 / 3", "The physics foundation – flex_common/")

    add_textbox(s, MARGIN_LEFT, CONTENT_TOP + Inches(0.06), Inches(12.2), Inches(0.3),
                "Seven small modules, ~1 300 lines. Each is one piece of textbook mechanics – "
                "this is what makes the model physics-based rather than a curve fit.",
                size=Pt(12.5), color=CARBON_GRAY)

    brand_table(
        s, MARGIN_LEFT, Inches(2.02), FULL_W,
        [Inches(1.62), Inches(0.72), Inches(9.945)],
        [
            ("baseline.py", "79",
             ("The reference part.", "One agreed set of numbers — R = 20 mm, wall 0.4 mm, "
              "ovalization 0.30 mm, 200/202 teeth — so every module means the same gearbox.")),
            ("ring.py", "170",
             ("The flexspline as a thin ring.", "Squash a thin steel ring into an ellipse; "
              "gives bending moment, shear and contact pressure in closed form.")),
            ("ellipse.py", "111",
             ("A real ellipse is not a pure cos 2θ.", "Decomposes the true shape into "
              "harmonics and adds the ring response to each. The correction term for ring.py.")),
            ("mesh.py", "305",
             ("Which teeth are touching, right now.", "From the wave-generator angle: which "
              "flexspline teeth line up with which spaces, and how each tooth moves — "
              "radial breathing, tangential slide, tilt. The heart of the kinematics.")),
            ("profile.py", "433",
             ("What shape the teeth must be.", "Synthesises the circular-spline flank as the "
              "conjugate envelope of the flexspline tooth, and the flank sliding speed — "
              "which feeds friction and wear.")),
            ("friction.py", "74",
             ("Why it takes torque to spin an unloaded drive.", "Without friction the no-load "
              "torque is exactly zero; this shows why, and what friction adds.")),
            ("teeth.py", "184",
             ("Tooth-root stress.", "Lewis cantilever bending plus Dolan–Broghamer fillet "
              "concentration — the fatigue surrogate.")),
        ],
        ["Module", "Lines", "In one sentence"],
        body_size=Pt(10.5), row_h=Inches(0.585),
    )

    add_textbox(s, MARGIN_LEFT, Inches(6.76), FULL_W, Inches(0.3),
                "A user of the twin never opens these — they are the foundation under the "
                "floorboards. But a skeptical reviewer goes straight here.",
                size=Pt(11), color=CARBON_GRAY)
    footer(s, 2)


def slide_dials_proof(prs):
    """Slide 3 — the dials, the proof, and how a user meets the package."""
    s = blank(prs)
    header(s, "Package anatomy – 3 / 3", "The dials, the proof, and how to start")

    # left — params.py
    add_textbox(s, MARGIN_LEFT, CONTENT_TOP + Inches(0.04), Inches(7.0), Inches(0.3),
                "The dials – params.py", size=Pt(14), bold=True, color=SCHAEFFLER_GREEN)
    add_textbox(s, MARGIN_LEFT, Inches(1.86), Inches(7.1), Inches(0.28),
                "Every adjustable quantity lives in one of four labelled boxes.",
                size=Pt(11), color=CARBON_GRAY)
    brand_table(
        s, MARGIN_LEFT, Inches(2.2), Inches(7.1),
        [Inches(1.82), Inches(2.08), Inches(3.2)],
        [
            ("FrictionParams", "How things rub",
             "boundary friction, Stribeck speed, WG bearing coefficient, grease churn"),
            ("ThermalParams", "How temperature changes things",
             "viscosity vs temperature, thermal expansion of clearances"),
            ("ToleranceParams", "Manufacturing imperfection",
             "backlash, tooth-to-tooth pitch scatter (Monte-Carlo)"),
            ("StiffnessChain", "How the drive twists",
             "series chain: cup + WG bearing + mesh + output bearing"),
        ],
        ["Dataclass", "What it holds", "Example knobs"],
        body_size=Pt(10), row_h=Inches(0.5),
    )
    add_textbox(s, MARGIN_LEFT, Inches(4.82), Inches(7.1), Inches(0.32),
                "A user adapting the twin to their gearbox edits here — never the equations. "
                "That is the contract.  Units: mm, N, N·mm, rad/s, °C.",
                size=Pt(11), color=CARBON_GRAY)

    # divider
    div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.96), CONTENT_TOP, Pt(1), Inches(4.35))
    div.fill.solid(); div.fill.fore_color.rgb = LIGHT_GRAY
    div.line.fill.background(); _no_shadow(div)

    # right — the proof
    x = Inches(8.32)
    add_textbox(s, x, CONTENT_TOP + Inches(0.04), Inches(4.5), Inches(0.3),
                "The proof – catalogue.py + tests/", size=Pt(14), bold=True,
                color=SCHAEFFLER_GREEN)
    keybox(s, x, Inches(1.84), Inches(4.49), Inches(1.42),
           "The headline claim",
           "Five independent catalogue quantities, four frame sizes, matched with only "
           "three global calibration dials and zero per-size fitting.")
    bullets(s, x, Inches(3.42), Inches(4.6), Inches(1.3), [
        "Schaeffler TPI 275, series RT1-H-CS, i = 100, +20 °C.",
        "Torsional stiffness uses no fitting at all — FE-anchored physical chain.",
        "Four tests keep the claim honest: if pytest passes, the README numbers are true.",
    ], size=Pt(10.5), spacing=Pt(5))

    # bottom — three depths
    add_textbox(s, MARGIN_LEFT, Inches(5.42), Inches(12.2), Inches(0.3),
                "How a user meets the package — most stop at the first depth",
                size=Pt(13), bold=True, color=SCHAEFFLER_GREEN)
    depths = [
        ("DEPTH 1", "\"just give me a number\"",
         "solve_point(30.0, 2000)  →  η = 77.9 %", "one line, ~30 s to first result"),
        ("DEPTH 2", "\"model my specific gearbox\"",
         "twin_for_size(17)", "or build a twin with your own parameter blocks"),
        ("DEPTH 3", "\"put it in my robot sim\"",
         "isaacsim6_actuator.py", "copy the class, call .step() each frame"),
    ]
    bx, bw, gap = MARGIN_LEFT, Inches(3.95), Inches(0.19)
    for i, (tag, what, code, note) in enumerate(depths):
        left = bx + i * (bw + gap)
        box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(5.8), bw, Inches(1.16))
        box.fill.solid(); box.fill.fore_color.rgb = GREEN_LIGHT
        box.line.fill.background(); _no_shadow(box)
        add_marker(s, left, Inches(5.8), Inches(1.16), width=Inches(0.07))
        add_textbox(s, left + Inches(0.18), Inches(5.87), Inches(3.6), Inches(0.2),
                    tag + "  ·  " + what, size=Pt(10), bold=True, color=SCHAEFFLER_GREEN)
        add_textbox(s, left + Inches(0.18), Inches(6.13), Inches(3.6), Inches(0.24),
                    code, size=Pt(10.5), bold=True, color=CARBON_GRAY, name=MONO_NAME)
        add_textbox(s, left + Inches(0.18), Inches(6.44), Inches(3.62), Inches(0.4),
                    note, size=Pt(9.5), color=CARBON_GRAY)

    # claim sits in the free space under the right column, clear of the footer
    s.shapes.add_picture(CLAIM_PATH, Inches(10.72), Inches(4.72), Inches(1.79), Inches(0.206))
    footer(s, 3)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide_overview(prs)
    slide_foundation(prs)
    slide_dials_proof(prs)
    prs.save(OUT)
    print("saved:", OUT, "slides:", len(prs.slides))


if __name__ == "__main__":
    main()
